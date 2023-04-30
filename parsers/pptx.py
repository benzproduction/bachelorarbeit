from pptx import Presentation
from typing import List, Tuple

def pptx_to_text(input_filename: str, slides_seperator: str = "\n\n") -> str:
    presentation = Presentation(input_filename)
    presentation_text = ""

    for slide in presentation.slides:
        
        slide_has_title = slide.shapes.title is not None

        for shape in slide.shapes:
            if not hasattr(shape, "text"):
                continue

            shape_text = f'\n{shape.text}'

            if slide_has_title and shape.text == slide.shapes.title.text:
                shape_text += ":"

            presentation_text += shape_text

        presentation_text += slides_seperator

    return presentation_text

def pptx_to_text_mapping(filename: str) -> List[Tuple[int, int, str]]:
    """
    Extracts the text from all slides of a PowerPoint (.pptx) file and returns it as a list of tuples.
    
    Args:
    - filename: str, the name or path of the PowerPoint file to read
    
    Returns:
    - A list of tuples, where each tuple has the following elements:
      * the slide number (starting from 0)
      * the offset in characters of the start of the slide text relative to the beginning of the presentation
      * the text extracted from the slide
      
    Raises:
    - FileNotFoundError if the given filename does not exist or cannot be opened
    - ValueError if the given filename does not correspond to a PowerPoint file
    
    Example usage:
    >>> extract_pptx_text('example.pptx')
    [(0, 0, 'Slide 1\nThis is some text on slide 1.\n\nThis is some more text on slide 1.'), 
     (1, 68, 'Slide 2\nThis is some text on slide 2.'), 
     (2, 96, 'Slide 3\nThis is some text on slide 3.\n\nThis is some more text on slide 3.'), 
     ...]
    """
    offset = 0
    slide_map = []
    
    prs = Presentation(filename)
    slides = prs.slides
    for slide_num, slide in enumerate(slides):
        shapes = slide.shapes
        text = '\n'.join([shape.text for shape in shapes if shape.has_text_frame])
        slide_map.append((slide_num, offset, text))
        offset += len(text)
        
    return slide_map